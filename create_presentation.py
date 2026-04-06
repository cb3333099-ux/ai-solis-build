"""
create_presentation.py

Generates a professional 10-slide PowerPoint presentation titled
"Contribution of Pāṇini to Linguistics" using python-pptx.

Output: panini_linguistics_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x1F, 0x3A, 0x6E)   # Deep navy blue
NAVY_DARK  = RGBColor(0x0D, 0x1F, 0x42)   # Darker navy for title bg
GOLD       = RGBColor(0xC8, 0xA0, 0x00)   # Academic gold accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)   # Slide background tint
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # Near-black body text
HIGHLIGHT  = RGBColor(0xC8, 0xA0, 0x00)   # Gold for highlighted terms

# ---------------------------------------------------------------------------
# Slide dimensions (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def set_slide_background(slide, color: RGBColor):
    """Fill the entire slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color: RGBColor = TEXT_DARK, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    """Add a simple text-box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def add_paragraph(text_frame, text, font_size, bold=False,
                  color: RGBColor = TEXT_DARK, align=PP_ALIGN.LEFT,
                  space_before=Pt(6), italic=False):
    """Append a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_mixed_paragraph(text_frame, segments, font_size,
                        align=PP_ALIGN.LEFT, space_before=Pt(6)):
    """
    Add a paragraph composed of multiple runs with different formatting.
    segments: list of (text, bold, color, italic)
    """
    p = text_frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    for seg_text, seg_bold, seg_color, seg_italic in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.size = font_size
        run.font.bold = seg_bold
        run.font.italic = seg_italic
        run.font.color.rgb = seg_color
    return p


def add_notes(slide, notes_text: str):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_divider(slide, top, color: RGBColor = GOLD):
    """Add a thin horizontal gold rule."""
    add_rect(slide, Inches(0.6), top, Inches(12.1), Inches(0.04),
             fill_color=color)


def add_image_placeholder(slide, left, top, width, height, label: str):
    """
    Add a styled rectangle that serves as an image placeholder,
    with a descriptive label centred inside it.
    """
    # Outer box
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.5)

    # Label text
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = NAVY


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def build_slide_1(prs: Presentation):
    """Title Slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY_DARK)

    # Top gold bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), GOLD)
    # Bottom gold bar
    add_rect(slide, Inches(0), Inches(7.32), SLIDE_W, Inches(0.18), GOLD)

    # Central white card
    add_rect(slide, Inches(1.2), Inches(1.4), Inches(10.93), Inches(4.7),
             RGBColor(0xFF, 0xFF, 0xFF))

    # Main title
    txBox = slide.shapes.add_textbox(
        Inches(1.4), Inches(1.7), Inches(10.5), Inches(2.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Contribution of Pāṇini to Linguistics"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = NAVY_DARK

    # Gold rule inside card
    add_rect(slide, Inches(3.5), Inches(3.6), Inches(6.33), Inches(0.05), GOLD)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(1.4), Inches(3.75), Inches(10.5), Inches(0.8)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Study of the Ancient Indian Grammarian's Enduring Legacy"
    run2.font.size = Pt(20)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x44, 0x55, 0x77)

    # Tertiary line
    txBox3 = slide.shapes.add_textbox(
        Inches(1.4), Inches(4.55), Inches(10.5), Inches(0.5)
    )
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Academic Presentation  ·  Sanskrit & Linguistic Studies"
    run3.font.size = Pt(13)
    run3.font.color.rgb = RGBColor(0x88, 0x88, 0xAA)

    # Decorative Sanskrit-style ornament text at bottom
    add_textbox(
        slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.5),
        "॥ अष्टाध्यायी ॥",
        Pt(16), bold=False, color=GOLD, align=PP_ALIGN.CENTER
    )

    add_notes(slide,
        "Welcome to this academic presentation on the profound contributions of Pāṇini "
        "to linguistics. Pāṇini, the ancient Indian grammarian who lived around the 4th "
        "century BCE, revolutionised our understanding of language through his monumental "
        "work, the Aṣṭādhyāyī. This presentation explores his key contributions that "
        "continue to influence modern linguistic science."
    )
    return slide


def build_slide_2(prs: Presentation):
    """Introduction — Background of Pāṇini."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    # Header band
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.85),
        "Introduction — Background of Pāṇini",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 2",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )

    # Gold divider below header
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    # Content text box
    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.6), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.space_before = Pt(0)
    run0 = p0.add_run()
    run0.text = (
        "Pāṇini (c. 4th century BCE) was a pioneering Sanskrit grammarian from "
        "Śalātura, a town in the northwestern region of ancient India (modern-day "
        "Pakistan). He is widely regarded as one of the greatest intellectuals of the "
        "ancient world and the father of descriptive and generative linguistics."
    )
    run0.font.size = Pt(14.5)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "Pāṇini lived during a remarkable period of intellectual flourishing in "
        "India. He was a student of the Vedic tradition and deeply acquainted with the "
        "oral literature of Sanskrit — a language of immense complexity and precision. "
        "His primary motivation was to preserve and codify the rules of classical "
        "Sanskrit in a rigorous, systematic manner.",
        Pt(14), color=TEXT_DARK, space_before=Pt(10)
    )

    add_paragraph(tf,
        "He composed his magnum opus, the Aṣṭādhyāyī, a grammar of Sanskrit consisting "
        "of nearly 3,959 tightly formulated rules (sūtras). These rules cover every "
        "aspect of language — from phonology and morphology to syntax — representing "
        "an extraordinary feat of intellectual engineering.",
        Pt(14), color=TEXT_DARK, space_before=Pt(10)
    )

    add_paragraph(tf,
        "His work predates modern formal linguistics by over two millennia, yet its "
        "precision, economy, and generative power are strikingly analogous to contemporary "
        "theoretical linguistics, mathematical logic, and even computer science.",
        Pt(14), color=TEXT_DARK, space_before=Pt(10)
    )

    # Image placeholder
    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.5),
        "📷  Pāṇini — Artistic depiction or\nstatue illustration\n(Replace with image)"
    )

    # Footer rule
    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "Pāṇini is believed to have lived around the 4th century BCE, though some "
        "scholars suggest an earlier date (6th–5th century BCE). His birthplace, Śalātura, "
        "was located near the Indus River. Remarkably, his grammar was composed in an "
        "oral tradition and later written down. Pāṇini's Aṣṭādhyāyī stands as the oldest "
        "known complete linguistic grammar of any language in the world. His work was so "
        "thorough that Sanskrit grammar remained essentially unchanged for centuries. "
        "Linguist Leonard Bloomfield described the Aṣṭādhyāyī as 'one of the greatest "
        "monuments of human intelligence.'"
    )
    return slide


def build_slide_3(prs: Presentation):
    """Aṣṭādhyāyī — His Main Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Aṣṭādhyāyī — The Grammar of Sanskrit",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 3",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    # Three key stat cards
    stats = [
        ("8", "Adhyāyas\n(Chapters)"),
        ("3,959", "Sūtras\n(Rules)"),
        ("~4th c. BCE", "Period of\nComposition"),
    ]
    card_w = Inches(2.6)
    card_h = Inches(1.2)
    card_top = Inches(1.45)
    for i, (num, label) in enumerate(stats):
        cx = Inches(0.55) + i * Inches(3.0)
        add_rect(slide, cx, card_top, card_w, card_h, NAVY)
        add_textbox(slide, cx, card_top + Inches(0.1), card_w, Inches(0.55),
                    num, Pt(26), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx, card_top + Inches(0.62), card_w, Inches(0.55),
                    label, Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

    # Main content
    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(2.85), Inches(8.5), Inches(4.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "The Aṣṭādhyāyī (meaning \"Eight Chapters\") is the foundational text of "
        "Sanskrit grammar and, by extension, a cornerstone of world linguistics. Written "
        "in an extremely concise sūtra style, each rule is crafted with maximum brevity "
        "and precision — often just a few syllables in length."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "Pāṇini organised the rules into eight chapters (adhyāyas), each subdivided "
        "into four pādas (quarters). The grammar covers phonology (śikṣā), morphology "
        "(pratyaya), syntax, and derivation. Each sūtra is interconnected with others "
        "through a sophisticated system of metalinguistic abbreviations called "
        "pratyāhāras.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "The economy of expression in the Aṣṭādhyāyī is remarkable: every element — "
        "every phoneme, suffix, and marker — serves a grammatical purpose. This principle "
        "of minimal redundancy (lāghava) is considered an intellectual breakthrough "
        "unmatched in the ancient world.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Modern scholars, including Ferdinand de Saussure and Noam Chomsky, have "
        "acknowledged Pāṇini's grammar as a precursor to formal language theory. "
        "The Aṣṭādhyāyī has been likened to a formal algorithm — a system capable of "
        "generating all grammatically correct Sanskrit sentences.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    # Image placeholder
    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.5),
        "📷  Page from a Sanskrit manuscript\nof the Aṣṭādhyāyī\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "The Aṣṭādhyāyī is a supreme example of scientific thinking in the ancient world. "
        "It is not merely a grammar book but a formal system for generating Sanskrit "
        "expressions. The text uses a special metalanguage with unique symbols (anubandhas) "
        "to make rules as concise as possible. For example, the rule 'ik yaṇ aci' in "
        "just three syllables encapsulates the sandhi (sound combination) rule for "
        "semi-vowel substitutions. The text has been transmitted intact for over 2,500 "
        "years, making it one of the most durable intellectual achievements in human history."
    )
    return slide


def build_slide_4(prs: Presentation):
    """Scientific and Logical Approach."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Scientific and Logical Approach to Grammar",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 4",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "Pāṇini's approach to grammar was fundamentally scientific. Rather than "
        "prescribing how language should be used (prescriptivism), he described the "
        "actual patterns of Sanskrit usage (descriptivism) with extraordinary accuracy "
        "— an approach not adopted systematically in Western linguistics until the 19th "
        "century."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "His method was deeply analytical: he identified the smallest meaningful units "
        "of language (morphemes), catalogued them, and formulated rules that predict "
        "their behaviour in every possible context. This mirrors modern scientific "
        "methodology — observation, hypothesis, and formalisation.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Pāṇini introduced a principle of rule economy called lāghava (brevity), which "
        "required that every sūtra should be as concise as possible without loss of "
        "generality. He also applied paribhāṣā — meta-rules that govern how other rules "
        "are to be interpreted — an idea central to modern formal grammars.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "His grammar treats language as a formal system with a finite set of rules "
        "capable of producing an infinite set of valid expressions. This logical "
        "architecture prefigures the Chomskyan notion of a finite grammar generating "
        "an infinite language — a key insight of 20th-century theoretical linguistics.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "The organisation of phonemes into sets using pratyāhāras (shorthand notation) "
        "demonstrates abstract algebraic thinking, grouping sounds by shared phonological "
        "properties — a method reminiscent of set theory in mathematics.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(5.0),
        "📷  Diagram of linguistic rule\nsystem / formal grammar\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "Pāṇini's scientific approach is striking because it is entirely empirical. "
        "He did not base his grammar on philosophical or religious authority but on the "
        "observable facts of Sanskrit usage. His lāghava principle is analogous to "
        "Occam's Razor in Western philosophy — prefer the simplest explanation. "
        "Pāṇini's paribhāṣās (meta-rules) are essentially rules about how to apply rules, "
        "which is a concept central to formal logic and computer science. "
        "His ability to reduce the entirety of Sanskrit grammar to fewer than 4,000 rules "
        "is considered one of the most impressive intellectual achievements in human history."
    )
    return slide


def build_slide_5(prs: Presentation):
    """Generative Grammar Concept."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Generative Grammar — A Pre-Modern Concept",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 5",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    # Quote box
    quote_box = slide.shapes.add_shape(1, Inches(0.55), Inches(1.42),
                                       Inches(12.2), Inches(1.0))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = NAVY
    quote_box.line.fill.background()
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    p_q = tf_q.paragraphs[0]
    p_q.alignment = PP_ALIGN.CENTER
    run_q = p_q.add_run()
    run_q.text = (
        "\"Pāṇini's grammar is the most complete generative grammar of any language "
        "ever written, and represents one of the greatest intellectual achievements "
        "of all time.\"  — Leonard Bloomfield"
    )
    run_q.font.size = Pt(13.5)
    run_q.font.italic = True
    run_q.font.color.rgb = GOLD

    # Main content
    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(2.6), Inches(8.5), Inches(4.65)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "Noam Chomsky introduced the theory of Generative Grammar in the 1950s, "
        "proposing that the human mind contains a set of innate rules capable of "
        "generating all grammatical sentences of a language. Remarkably, Pāṇini "
        "had articulated a strikingly similar framework over two thousand years earlier."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "Pāṇini's Aṣṭādhyāyī functions precisely as a generative grammar: it provides "
        "a finite set of rules (sūtras) which, when applied recursively and in a specific "
        "order, can generate the complete set of well-formed Sanskrit expressions. "
        "This includes derivational processes, inflectional paradigms, and sandhi rules.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "The concept of a 'derivation' in Pāṇinian grammar — starting from a base "
        "(dhātu or prātipadika) and progressively adding affixes according to rules — "
        "directly parallels the transformational derivations in Chomskyan syntax. "
        "Both systems are rule-governed, recursive, and productive.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Scholars such as Paul Kiparsky and John Staal have demonstrated formal "
        "equivalences between Pāṇinian grammar and modern transformational-generative "
        "grammar, establishing Pāṇini as the world's first generative linguist.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(5.0),
        "📷  Chomsky hierarchy / Generative\ngrammar tree diagram\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "The idea of generative grammar — that a finite set of rules can produce an "
        "infinite number of valid sentences — is central to modern linguistics. "
        "Chomsky is credited with formalising this in the 20th century, but historians "
        "of linguistics recognise Pāṇini as having achieved the same insight. "
        "The Pāṇinian derivational process begins with a semantic intention (vivakṣā) "
        "and passes through phonological, morphological, and syntactic rules to arrive "
        "at the final linguistic form. This pipeline is architecturally identical to "
        "the derivation trees used in generative syntax today."
    )
    return slide


def build_slide_6(prs: Presentation):
    """Use of Meta-language."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Use of Meta-language in the Aṣṭādhyāyī",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 6",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "One of Pāṇini's most revolutionary contributions was his invention of a "
        "formal meta-language — a language used to talk about language itself. "
        "His meta-language employed special symbols, abbreviated codes, and conventions "
        "that allowed him to express complex grammatical rules with extraordinary brevity."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "The most important device in this meta-language is the pratyāhāra — a system "
        "of phonological abbreviations. Using the sounds listed in the Śivasūtras "
        "(a pre-pended list of Sanskrit phonemes), Pāṇini coined short codes like 'aṇ' "
        "or 'ac' to refer to entire classes of sounds, reducing verbose descriptions "
        "to single syllables.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Another device is the anubandha (IT marker) — a technical indicator attached "
        "to grammatical elements to signal their functional role. These markers are "
        "mnemonic and operational: they tell the grammar how to apply a rule, but are "
        "not part of the output form. This resembles the variable symbols in mathematical "
        "equations or the dummy symbols in computer programming.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Pāṇini's meta-language thus anticipates modern formal language theory and the "
        "use of symbolic notation in mathematics, logic, and computer science. "
        "Linguist Ferdinand de Saussure acknowledged the influence of Pāṇini's "
        "symbolic approach on his own structuralist framework.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "The Śivasūtras, a 14-line phoneme inventory, can be seen as the first known "
        "use of an artificial formal notation system for linguistic description — "
        "a feat that predates Boolean algebra and formal logic by centuries.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.5),
        "📷  Illustration of pratyāhāra /\nŚivasūtra notation\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "The Śivasūtras list all Sanskrit phonemes in a specific order that allows "
        "classes of sounds to be referred to by combining the first and last element "
        "(e.g., 'a-ṇ' covers all vowels). This device allows Pāṇini to write 'ik yaṇ aci' "
        "to mean 'the sounds i, u, ṛ, ḷ are replaced by their semi-vowel counterparts "
        "when followed by a vowel.' This kind of compact formal notation was not used "
        "in Western linguistics until the 20th century. The concept of IT markers "
        "(anubandhas) is especially notable: they function like dummy variables in "
        "programming — they serve a computational purpose but are invisible in the output."
    )
    return slide


def build_slide_7(prs: Presentation):
    """Phonetics and Sound Classification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Phonetics and Sound Classification",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 7",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "Pāṇini's work on phonetics (śikṣā) represents one of the earliest and most "
        "precise accounts of the human sound system. His classification of Sanskrit "
        "phonemes was based on rigorous articulatory phonetics — identifying where and "
        "how sounds are produced in the vocal tract."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "He classified sounds according to their place of articulation: velars (kaṇṭhya), "
        "palatals (tālavya), retroflexes (mūrdhanya), dentals (dantya), and labials "
        "(oṣṭhya). He further distinguished sounds by manner of articulation — stops, "
        "fricatives, nasals, approximants — with remarkable accuracy.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "His treatment of vowels (svara) and consonants (vyañjana), their lengths "
        "(hrasva, dīrgha, pluta), and their tonal qualities (udātta, anudātta, svarita) "
        "provides a complete phonological analysis of Sanskrit. This degree of phonetic "
        "description was not matched in Europe until the 18th–19th centuries.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "His phonological rules — particularly the sandhi rules describing how sounds "
        "change at word boundaries and morpheme junctions — represent a systematic "
        "study of allophonic variation and phonological conditioning, core concepts in "
        "modern phonology.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "The International Phonetic Alphabet (IPA) developed in the 19th century shares "
        "several classification principles with Pāṇinian phonetics, and some scholars "
        "argue that Indian phonological science directly influenced early comparative "
        "linguists such as William Jones and Franz Bopp.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.8),
        "📷  Sanskrit phoneme chart /\narticulatory phonetics diagram\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "The Śikṣā texts (phonetic treatises) associated with the Vedic tradition were "
        "an important background for Pāṇini's phonological work. His system recognises "
        "over 60 distinct phonemes in Sanskrit, classified by place and manner of "
        "articulation — far more precise than anything found in ancient Greek or Latin "
        "grammatical traditions. Sandhi rules — rules governing sound changes at "
        "morpheme and word boundaries — cover both internal sandhi (within a word) and "
        "external sandhi (between words). These rules are phonologically conditioned "
        "and context-sensitive, paralleling modern phonological rule notation."
    )
    return slide


def build_slide_8(prs: Presentation):
    """Morphology (Word Formation)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Morphology — Word Formation and Structure",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 8",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "Morphology — the study of word formation — is arguably the most thoroughly "
        "developed aspect of Pāṇini's grammar. He provided a comprehensive account of "
        "how Sanskrit words are built from roots (dhātus), suffixes (pratyayas), and "
        "augments, covering both inflectional and derivational morphology."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "Pāṇini catalogued over 2,000 verbal roots in the Dhātupāṭha, classifying "
        "them into ten conjugational classes (gaṇas) based on the pattern of their "
        "present tense stem formation. For each root, he could predict all possible "
        "inflected and derived forms through the application of ordered rules.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "His treatment of nominal morphology (nāmaprakaraṇa) covers all eight cases "
        "(vibhaktis) — nominative, accusative, instrumental, dative, ablative, "
        "genitive, locative, and vocative — across all genders and declensional "
        "classes. This systematic account of inflectional morphology is still used "
        "as the basis for Sanskrit pedagogy today.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Pāṇini's morphological framework introduced the notion of zero morpheme "
        "(lopa) — a rule that a grammatical element may be phonetically null but "
        "grammatically present. This concept was independently rediscovered in modern "
        "morphology and is standard in contemporary linguistic analysis.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "His concept of kāraka — the semantic roles (agent, patient, instrument, etc.) "
        "that nouns play in a sentence — anticipates modern case grammar and thematic "
        "role theory, providing a bridge between morphology and syntax.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.8),
        "📷  Sanskrit morphology chart /\nword formation diagram\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "Sanskrit morphology is extraordinarily complex — a single verbal root can "
        "theoretically generate hundreds of derived forms. Pāṇini's system handles "
        "this complexity through a derivational pipeline: root + class marker + "
        "tense/mood suffix + personal ending, with sandhi rules applied at each stage. "
        "The kāraka theory is particularly significant for linguistics. It differentiates "
        "between surface case (vibhakti) and semantic role (kāraka), a distinction "
        "that became central to Fillmore's Case Grammar (1968) in modern linguistics. "
        "The concept of zero morpheme (lopa) — that grammatical information can be "
        "encoded by the absence of a phoneme — is used in analyses of English "
        "plurals (sheep → sheep-Ø) and past tenses."
    )
    return slide


def build_slide_9(prs: Presentation):
    """Rule Ordering and Influence on Modern Linguistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.25), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.85),
        "Rule Ordering and Influence on Modern Linguistics",
        Pt(28), bold=True, color=WHITE
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.45),
        "Slide 9",
        Pt(11), color=RGBColor(0xAA, 0xBB, 0xCC)
    )
    add_rect(slide, Inches(0), Inches(1.25), SLIDE_W, Inches(0.055), GOLD)

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.45), Inches(8.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = (
        "A defining feature of Pāṇini's grammar is the strict ordering of rules. "
        "When two or more rules are applicable to the same input, Pāṇini specifies "
        "precedence through meta-rules such as 'antaraṅga' (inner rule) taking "
        "priority over 'bahiraṅga' (outer rule), or 'nitya' (obligatory) over "
        "'anitya' (optional). This is called 'pāṇinian rule ordering'."
    )
    run0.font.size = Pt(14)
    run0.font.color.rgb = TEXT_DARK

    add_paragraph(tf,
        "Rule ordering is central to modern generative phonology and morphology. "
        "Chomsky and Halle's The Sound Pattern of English (1968) — one of the most "
        "influential works in modern linguistics — explicitly employs ordered rules, "
        "a concept directly traceable to Pāṇinian grammar.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Pāṇini's influence on the development of modern linguistics is well "
        "documented. When European scholars encountered the Aṣṭādhyāyī in the "
        "18th–19th centuries — through William Jones's discovery of Sanskrit's "
        "kinship with Greek and Latin — it sparked the entire field of comparative "
        "and historical linguistics. Scholars like Franz Bopp studied Sanskrit "
        "grammar intensively, shaping the early development of Indo-European studies.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_paragraph(tf,
        "Beyond linguistics, Pāṇini's ideas have influenced computer science. "
        "Backus–Naur Form (BNF), used to define programming language syntax, "
        "is structurally similar to Pāṇini's rule notation. The Indian computer "
        "scientist Panini Backus Form hypothesis even suggests that BNF was "
        "independently converged upon Pāṇinian principles.",
        Pt(14), color=TEXT_DARK, space_before=Pt(8)
    )

    add_image_placeholder(
        slide, Inches(9.35), Inches(1.45), Inches(3.7), Inches(4.8),
        "📷  Timeline of Pāṇini's influence\non modern linguistics\n(Replace with image)"
    )

    add_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.06), NAVY)

    add_notes(slide,
        "The concept of rule ordering is subtle but crucial. Without it, many grammars "
        "become ambiguous or over-generate incorrect forms. Pāṇini's system resolves "
        "these conflicts through a hierarchical system of meta-principles. "
        "In modern phonology, 'counter-bleeding' and 'counter-feeding' interactions "
        "between rules are still debated — these are exactly the kinds of interactions "
        "Pāṇini grappled with. William Jones's 1786 address to the Asiatic Society, "
        "noting the systematic correspondences between Sanskrit, Greek, and Latin, "
        "is traditionally taken as the birth of comparative linguistics. "
        "This discovery was only possible because Sanskrit's grammar had been preserved "
        "in such rigorous detail by Pāṇini."
    )
    return slide


def build_slide_10(prs: Presentation):
    """Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY_DARK)

    # Top gold bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), GOLD)
    # Bottom gold bar
    add_rect(slide, Inches(0), Inches(7.32), SLIDE_W, Inches(0.18), GOLD)

    # Heading
    add_textbox(
        slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
        "Conclusion",
        Pt(32), bold=True, color=GOLD, align=PP_ALIGN.LEFT
    )

    # Thin white rule
    add_rect(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.04), WHITE)

    # Content text box
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.1), Inches(9.2), Inches(6.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.space_before = Pt(0)
    run0 = p0.add_run()
    run0.text = (
        "Pāṇini's contribution to linguistics is monumental, timeless, and "
        "far-reaching. Living more than 2,500 years ago, he produced a work "
        "of such intellectual depth and formal precision that it continues to "
        "influence linguistic science, computer science, and philosophy to this day."
    )
    run0.font.size = Pt(14.5)
    run0.font.color.rgb = WHITE

    add_paragraph(tf,
        "His Aṣṭādhyāyī stands as the world's first complete formal grammar — "
        "a rigorous, generative, and descriptively adequate account of an entire "
        "natural language. Its design principles — economy, recursion, ordered "
        "rules, meta-language — are the same principles that underpin modern "
        "theoretical linguistics and formal language theory.",
        Pt(14), color=RGBColor(0xCC, 0xD5, 0xE0), space_before=Pt(10)
    )

    add_paragraph(tf,
        "Pāṇini's work on phonetics, morphology, syntax, and meta-language places "
        "him centuries, if not millennia, ahead of his time. Scholars continue to "
        "find new connections between his system and modern linguistic theories, "
        "demonstrating the inexhaustible depth of his intellectual vision.",
        Pt(14), color=RGBColor(0xCC, 0xD5, 0xE0), space_before=Pt(10)
    )

    add_paragraph(tf,
        "In summary, Pāṇini was not merely a grammarian — he was the world's "
        "first formal linguist, and his legacy is a testament to the power of "
        "rigorous, systematic thought in unlocking the secrets of human language.",
        Pt(14.5), bold=True, color=GOLD, space_before=Pt(14)
    )

    # Key takeaways panel
    box_right = slide.shapes.add_shape(
        1, Inches(10.0), Inches(1.1), Inches(3.0), Inches(5.8)
    )
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x55)
    box_right.line.fill.background()
    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.alignment = PP_ALIGN.CENTER
    run_r = p_r.add_run()
    run_r.text = "Key Takeaways"
    run_r.font.size = Pt(14)
    run_r.font.bold = True
    run_r.font.color.rgb = GOLD

    takeaways = [
        "First formal grammar\nof any world language",
        "Precursor to generative\ngrammar theory",
        "Inventor of formal\nmeta-language",
        "Pioneer of scientific\nphonology",
        "Lasting influence on\nlinguistics & CS",
    ]
    for i, t in enumerate(takeaways):
        txB = slide.shapes.add_textbox(
            Inches(10.05), Inches(1.75) + i * Inches(0.9),
            Inches(2.9), Inches(0.85)
        )
        tf_t = txB.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.alignment = PP_ALIGN.CENTER
        r_t = p_t.add_run()
        r_t.text = f"▸  {t}"
        r_t.font.size = Pt(11)
        r_t.font.color.rgb = WHITE

    # Decorative Sanskrit text
    add_textbox(
        slide, Inches(0.6), Inches(7.0), Inches(9), Inches(0.3),
        "॥ संस्कृत — The Language of Perfection ॥",
        Pt(12), color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.LEFT
    )

    add_notes(slide,
        "In conclusion, Pāṇini's contributions to linguistics are unique in world "
        "intellectual history. No other ancient scholar produced a work of equivalent "
        "rigour, generality, and lasting relevance. The Aṣṭādhyāyī has been in "
        "continuous use as the authoritative grammar of Sanskrit for over 2,500 years. "
        "Its influence extends to comparative linguistics (William Jones, Franz Bopp), "
        "formal language theory (Chomsky hierarchy), computer science (BNF notation), "
        "and philosophy of language (Bhartrhari's continuation). Pāṇini reminds us "
        "that great science transcends the boundaries of time and culture. "
        "Thank the audience for their attention and invite questions."
    )
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slide 1  — Title Slide")
    build_slide_1(prs)

    print("Building slide 2  — Introduction")
    build_slide_2(prs)

    print("Building slide 3  — Aṣṭādhyāyī")
    build_slide_3(prs)

    print("Building slide 4  — Scientific Approach")
    build_slide_4(prs)

    print("Building slide 5  — Generative Grammar")
    build_slide_5(prs)

    print("Building slide 6  — Meta-language")
    build_slide_6(prs)

    print("Building slide 7  — Phonetics")
    build_slide_7(prs)

    print("Building slide 8  — Morphology")
    build_slide_8(prs)

    print("Building slide 9  — Rule Ordering & Influence")
    build_slide_9(prs)

    print("Building slide 10 — Conclusion")
    build_slide_10(prs)

    output_file = "panini_linguistics_presentation.pptx"
    prs.save(output_file)
    print(f"\n✅  Presentation saved as '{output_file}'")
    print(f"   Slides: {len(prs.slides)}")
    print("   Open the file in Microsoft PowerPoint or LibreOffice Impress.")


if __name__ == "__main__":
    create_presentation()
